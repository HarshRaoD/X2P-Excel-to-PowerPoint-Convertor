# Basic Libraries
import numpy as np
import pandas as pd
import seaborn as sb
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
import tkinter as tk
from tkinter import filedialog
from tkinter import messagebox
import os
from datetime import datetime

window = tk.Tk()
window.title("X2P - Excel to Powerpoint Convertor")

#Space to declare global variables starts
presentationTitle = tk.StringVar()
presentationSubtitle = tk.StringVar()
colour_list = [["blue", "cyan", "green", "lime", "red", "magenta", "violet", "purple","yellow","orange"],
               ["cyan", "blue", "lime", "green", "magenta", "red", "m", "violet", "orange", "yellow"],
               ["Blues_r", "Blues", "Greens_r", "Greens", "Reds_r", "Purples_r", "Purples", "Purples_r", "autumn_r", 'autumn']]
# Row 0 is the primary colour, Row 1 is secondary Colour(for line graph), Row 2 is colour scheme for pie plot

Histogram = [["Select Column",""],["Select Column",""]] # Col 0 is column(data), Col 1 is graph title
Violin = [["Select Column","",0],["Select Column","",0]] # Col 0 is column(data), Col 1 is graph title, Col 2 is graph colour
Pie = [["Select Column","",1,0], ["Select Column","",1,0]] # Col 2 is for 1 - Use Others, 0 is no Use others, Col 3 is colour
Count = [["Select Column","",0],["Select Column","",0]] # Col 0 is column(data), Col 1 is graph title, Col 2 for colour
Bar = [["Select Column","Select Column","",0,0,"Select Column",0],["Select Column","Select Column","",0,0,"Select Column",0]] # Col 0 is cat data, Col1 num data, Col2 Title, Col3 (sum(0)/mean(1)), Col4 (invert axis 0/1), Col5 Hue, Col 6 colour
Line = [["Select Column","Select Column","",0,"Select Column",0,0],["Select Column","Select Column","",0,"Select Column",0,0]] # Col0 is X axis, Col1 Y axis, Col2 Title, Col3 Line design, Col4 2nd Y axis, Col5 2nd Y axis line Design, Col6 colour
#Space to declare global variables ends

def LinearSearch_index(l, a): # list l & element a
    '''returns index if a is in l or -1 if not, treats Select Column as an element by default'''
    if(a == 'Select Column'):
        return 1
    i = 0
    while(i < len(l)):
        if(l[i] == a):
            return i
        i += 1
    return -1

def LinearSearch_YesNo(l, a): # list l & element a
    '''returns 1 if a in l or 0 if not'''
    if(a == 'Select Column'):
        return 1
    i = 0
    while(i < len(l)):
        if(l[i] == a):
            return 1
        i += 1
    return 0
    
def pie_find_max_slices(nos): 
    # nos is a list of integers
    ''' To find out the maximum no of slices such that 'Others' < 10 %'''
    if(len(nos) <= 5):
        return 5
    i = len(nos) - 1
    sum_nos = 0
    for c in nos:
        sum_nos += c
    sum_nos1 = 0
    while(i >= 0):
       sum_nos1 += nos[i]
       if(sum_nos1 >= (sum_nos*10/100)): # Others cannot be >= 10%
           return (i + 2)
       i -= 1
    return 5

def check_column(col, dt): # Where col -> column, dt -> datatype expected in template
    ''''This function is for checking if a column is in the requiered datatype and converting it if not.
        It returns 1 if the conversion is successful or not needed and returns 0 if there is a problem.'''
    # dt can be 0->Not a column, 1->Numerical, 2->Categorical, 3->Timeseries
    if(col == 'Select Column'):
        return 1
    if(LinearSearch_YesNo(numerical_data, col) == 1):
        if(dt == 1):
            return 1 # Current and expected datatype match
        else:
            result = convert_datatype_2(col, dt)
    elif(LinearSearch_YesNo(categorical_data, col) == 1):
        if(dt == 2):
            return 1 # Current and expected datatype match
        else:
            result = convert_datatype_2(col, dt)
    elif(LinearSearch_YesNo(timeseries_data, col) == 1):
        if(dt == 3):
            return 1 # Current and expected datatype match
        else:
            result = convert_datatype_2(col, dt)
    else:
        return 0 # Column does not exist in data file
    #print("Result for check_column(",col,", ",dt,") is ",result)
    return result

def convert_datatype_2(col, dt):
    '''This Function converts the datatype of column (col) into that of datatype (dt)
        Returns 1 if successful and 0 if there is an error.
        This differs from convert_datatype() as that is interactive and used for
        explicit convertions while this is used for Implicit Conversions'''
    if(dt == 1): # if number
        # Now convert to numerical
        try:
            df[col] = pd.to_numeric(df[col])
        except ValueError as err:
            # if some error ask user if he wants to forcibly convert
            forcibly = messagebox.askokcancel("Unable to convert","There was a problem while loading your template. \n Press Ok to continue loading the template forcibly \n(This may result in loss of data)")
            if(forcibly == 1):
                df[col] = pd.to_numeric(df[col], errors='coerce')
                # Removing the column from its old type list
                if(LinearSearch_YesNo(categorical_data, col) == 1):
                    categorical_data.remove(col)
                else:
                    timeseries_data.remove(col)
                # Add the same column to new type list
                numerical_data.append(col)
            else:
                return 0
        else:
            '''If try block is executed'''
            # Removing the column from its old type list
            if(LinearSearch_YesNo(categorical_data, col) == 1):
                categorical_data.remove(col)
            else:
                timeseries_data.remove(col)
            # Add the same column to new type list
            numerical_data.append(col) 
    elif(dt == 3):
        try:
            df[col] = df[col].astype('datetime64[ns]')
        except ValueError as err:
            return 0
        else:
            '''If try block is executed'''
            # Removing the column from its old type list
            if(LinearSearch_YesNo(categorical_data, col) == 1):
                categorical_data.remove(col)
            else:
                numerical_data.remove(col)
            # Add the same column to new type list
            timeseries_data.append(col) 
    elif(dt == 2):
        df[col] = df[col].astype('object')
        # Removing the column from its old type list
        if(LinearSearch_YesNo(timeseries_data, col) == 1):
            timeseries_data.remove(col)
        else:
            numerical_data.remove(col)
        # Add the same column to new type list
        categorical_data.append(col) 

    return 1
    

def createTemplate(no):
    '''Çreates a text File to store the values of all the global lists'''
    if(no == 1): #Template 1
        with open("AppFinal_Template1.txt", "w") as f:
            # Histogram
            i = 0
            while(i < 2):
                j = 0
                while(j < 2):
                    print(Histogram[i][j], file = f)
                    j += 1
                i += 1
            # Violin
            i = 0
            while(i < 2):
                j = 0
                while(j < 3):
                    print(Violin[i][j], file = f)
                    j += 1
                i += 1
            # Pie
            i = 0
            while(i < 2):
                j = 0
                while(j < 4):
                    print(Pie[i][j], file = f)
                    j += 1
                i += 1
            # Count
            i = 0
            while(i < 2):
                j = 0
                while(j < 3):
                    print(Count[i][j], file = f)
                    j += 1
                i += 1
            # Bar
            i = 0
            while(i < 2):
                j = 0
                while(j < 7):
                    print(Bar[i][j], file = f)
                    j += 1
                i += 1
            # Line
            i = 0
            while(i < 2):
                j = 0
                while(j < 7):
                    print(Line[i][j], file = f)
                    j += 1
                i += 1
    elif(no == 2): #Template 2
        with open("AppFinal_Template2.txt", "w") as f:
            # Histogram
            i = 0
            while(i < 2):
                j = 0
                while(j < 2):
                    print(Histogram[i][j], file = f)
                    j += 1
                i += 1
            # Violin
            i = 0
            while(i < 2):
                j = 0
                while(j < 3):
                    print(Violin[i][j], file = f)
                    j += 1
                i += 1
            # Pie
            i = 0
            while(i < 2):
                j = 0
                while(j < 4):
                    print(Pie[i][j], file = f)
                    j += 1
                i += 1
            # Count
            i = 0
            while(i < 2):
                j = 0
                while(j < 3):
                    print(Count[i][j], file = f)
                    j += 1
                i += 1
            # Bar
            i = 0
            while(i < 2):
                j = 0
                while(j < 7):
                    print(Bar[i][j], file = f)
                    j += 1
                i += 1
            # Line
            i = 0
            while(i < 2):
                j = 0
                while(j < 7):
                    print(Line[i][j], file = f)
                    j += 1
                i += 1
    elif(no == 3): #Template 3
        with open("AppFinal_Template3.txt", "w") as f:
            # Histogram
            i = 0
            while(i < 2):
                j = 0
                while(j < 2):
                    print(Histogram[i][j], file = f)
                    j += 1
                i += 1
            # Violin
            i = 0
            while(i < 2):
                j = 0
                while(j < 3):
                    print(Violin[i][j], file = f)
                    j += 1
                i += 1
            # Pie
            i = 0
            while(i < 2):
                j = 0
                while(j < 4):
                    print(Pie[i][j], file = f)
                    j += 1
                i += 1
            # Count
            i = 0
            while(i < 2):
                j = 0
                while(j < 3):
                    print(Count[i][j], file = f)
                    j += 1
                i += 1
            # Bar
            i = 0
            while(i < 2):
                j = 0
                while(j < 7):
                    print(Bar[i][j], file = f)
                    j += 1
                i += 1
            # Line
            i = 0
            while(i < 2):
                j = 0
                while(j < 7):
                    print(Line[i][j], file = f)
                    j += 1
                i += 1

def load_template1(event):
    '''Reads the .txt for template 1 and loads it into the global arrays'''
    listTxt2 = []
    try:
        with open("AppFinal_Template1.txt") as f:
            for l in f.readlines(): # Goes line by line
                listTxt2.append(l.replace('\n',''))
    except FileNotFoundError:
        messagebox.showwarning("No Template","Template 1 hasn't been saved yet or has been deleted.")
        return
    r = 0
    # Histogram
    i = 0
    while(i < 2):
        j = 0
        while(j < 2):
            Histogram[i][j] = listTxt2[r]
            r += 1
            j += 1
        if(check_column(Histogram[i][0], 1) == 0):
            Histogram[i][0] = 'Select Column'
            messagebox.showerror("Incompatible Template","Template 1 is incompatible with the given data file. Please select the graphs manually")
            return
        i += 1
    
    # Violin (r = 3)
    i = 0
    while(i < 2):
        j = 0
        while(j < 3):
            if(j == 2):
                Violin[i][j] = int(listTxt2[r])
            else:
                Violin[i][j] = listTxt2[r]
            r += 1
            j += 1
        if(check_column(Violin[i][0], 1) == 0):
            Violin[i][0] = 'Select Column'
            messagebox.showerror("Incompatible Template","Template 1 is incompatible with the given data file. Please select the graphs manually")
            return
        i += 1
    # Pie (r = 9)
    i = 0
    while(i < 2):
        j = 0
        while(j < 4):
            if(j == 2 or j == 3):
                Pie[i][j] = int(listTxt2[r])
            else:
                Pie[i][j] = listTxt2[r]
            r += 1
            j += 1
        if(check_column(Pie[i][0], 2) == 0):
            Pie[i][0] = 'Select Column'
            messagebox.showerror("Incompatible Template","Template 1 is incompatible with the given data file. Please select the graphs manually")
            return
        i += 1
    # Count (r = 17)
    i = 0
    while(i < 2):
        j = 0
        while(j < 3):
            if(j == 2):
                Count[i][j] = int(listTxt2[r])
            else:
                Count[i][j] = listTxt2[r]
            r += 1
            j += 1
        if(check_column(Count[i][0], 2) == 0):
            Count[i][0] = 'Select Column'
            messagebox.showerror("Incompatible Template","Template 1 is incompatible with the given data file. Please select the graphs manually")
            return
        i += 1
    # Bar (r = 23)
    i = 0
    while(i < 2):
        j = 0
        while(j < 7):
            if(j == 3 or j == 4 or j == 6):
                Bar[i][j] = int(listTxt2[r])
            else:
                Bar[i][j] = listTxt2[r]
            r += 1
            j += 1
        if(check_column(Bar[i][0], 2) == 0 or check_column(Bar[i][1], 1) == 0 or check_column(Bar[i][5], 2) == 0):
            Bar[i][0] = 'Select Column'
            Bar[i][1] = 'Select Column'
            messagebox.showerror("Incompatible Template","Template 1 is incompatible with the given data file. Please select the graphs manually")
            return
        i += 1
    # Line (r = 37)
    i = 0
    while(i < 2):
        j = 0
        while(j < 7):
            if(j == 3 or j == 5 or j == 6):
                Line[i][j] = int(listTxt2[r])
            else:
                Line[i][j] = listTxt2[r]
            r += 1
            j += 1
        if(LinearSearch_YesNo(numerical_data + timeseries_data, Line[i][0]) == 0 or LinearSearch_YesNo(numerical_data + timeseries_data, Line[i][1]) == 0 or LinearSearch_YesNo(numerical_data + timeseries_data, Line[i][4]) == 0):
            ''''For Line Graph only: If datatype is converted from Categorical to others it cannot be stored in the template and manual entry needed'''
            print(Line)
            Line[i][0] = 'Select Column'
            Line[i][1] = 'Select Column'
            Line[i][4] = 'Select Column'
            messagebox.showerror("Incompatible Template","Template 1 is incompatible with the given data file. Please select the graphs manually")
            return
        i += 1
        
    messagebox.showinfo("Template Loaded","Template 1 has been loaded successfully.")
    return

def load_template2(event):
    '''Reads the .txt for template 2 and loads it into the global arrays'''
    listTxt2 = []
    try:
        with open("AppFinal_Template2.txt") as f:
            for l in f.readlines(): # Goes line by line
                listTxt2.append(l.replace('\n',''))
    except FileNotFoundError:
        messagebox.showwarning("No Template","Template 2 hasn't been saved yet or has been deleted.")
        return
    r = 0
    # Histogram
    i = 0
    while(i < 2):
        j = 0
        while(j < 2):
            Histogram[i][j] = listTxt2[r]
            r += 1
            j += 1
        if(check_column(Histogram[i][0], 1) == 0):
            Histogram[i][0] = 'Select Column'
            messagebox.showerror("Incompatible Template","Template 2 is incompatible with the given data file. Please select the graphs manually")
            return
        i += 1
    
    # Violin (r = 3)
    i = 0
    while(i < 2):
        j = 0
        while(j < 3):
            if(j == 2):
                Violin[i][j] = int(listTxt2[r])
            else:
                Violin[i][j] = listTxt2[r]
            r += 1
            j += 1
        if(check_column(Violin[i][0], 1) == 0):
            Violin[i][0] = 'Select Column'
            messagebox.showerror("Incompatible Template","Template 2 is incompatible with the given data file. Please select the graphs manually")
            return
        i += 1
    # Pie (r = 9)
    i = 0
    while(i < 2):
        j = 0
        while(j < 4):
            if(j == 2 or j == 3):
                Pie[i][j] = int(listTxt2[r])
            else:
                Pie[i][j] = listTxt2[r]
            r += 1
            j += 1
        if(check_column(Pie[i][0], 2) == 0):
            Pie[i][0] = 'Select Column'
            messagebox.showerror("Incompatible Template","Template 2 is incompatible with the given data file. Please select the graphs manually")
            return
        i += 1
    # Count (r = 17)
    i = 0
    while(i < 2):
        j = 0
        while(j < 3):
            if(j == 2):
                Count[i][j] = int(listTxt2[r])
            else:
                Count[i][j] = listTxt2[r]
            r += 1
            j += 1
        if(check_column(Count[i][0], 2) == 0):
            Count[i][0] = 'Select Column'
            messagebox.showerror("Incompatible Template","Template 2 is incompatible with the given data file. Please select the graphs manually")
            return
        i += 1
    # Bar (r = 23)
    i = 0
    while(i < 2):
        j = 0
        while(j < 7):
            if(j == 3 or j == 4 or j == 6):
                Bar[i][j] = int(listTxt2[r])
            else:
                Bar[i][j] = listTxt2[r]
            r += 1
            j += 1
        if(check_column(Bar[i][0], 2) == 0 or check_column(Bar[i][1], 1) == 0 or check_column(Bar[i][5], 2) == 0):
            Bar[i][0] = 'Select Column'
            Bar[i][1] = 'Select Column'
            messagebox.showerror("Incompatible Template","Template 2 is incompatible with the given data file. Please select the graphs manually")
            return
        i += 1
    # Line (r = 37)
    i = 0
    while(i < 2):
        j = 0
        while(j < 7):
            if(j == 3 or j == 5 or j == 6):
                Line[i][j] = int(listTxt2[r])
            else:
                Line[i][j] = listTxt2[r]
            r += 1
            j += 1
        if(LinearSearch_YesNo(numerical_data + timeseries_data, Line[i][0]) == 0 or LinearSearch_YesNo(numerical_data + timeseries_data, Line[i][1]) == 0 or LinearSearch_YesNo(numerical_data + timeseries_data, Line[i][4]) == 0):
            ''''For Line Graph only: If datatype is converted from Categorical to others it cannot be stored in the template and manual entry needed'''
            print(Line)
            Line[i][0] = 'Select Column'
            Line[i][1] = 'Select Column'
            Line[i][4] = 'Select Column'
            messagebox.showerror("Incompatible Template","Template 2 is incompatible with the given data file. Please select the graphs manually")
            return
        i += 1
        
    messagebox.showinfo("Template Loaded","Template 2 has been loaded successfully.")
    return
            
def load_template3(event):
    '''Reads the .txt for template 3 and loads it into the global arrays'''
    listTxt2 = []
    try:
        with open("AppFinal_Template3.txt") as f:
            for l in f.readlines(): # Goes line by line
                listTxt2.append(l.replace('\n',''))
    except FileNotFoundError:
        messagebox.showwarning("No Template","Template 3 hasn't been saved yet or has been deleted.")
        return
    r = 0
    # Histogram
    i = 0
    while(i < 2):
        j = 0
        while(j < 2):
            Histogram[i][j] = listTxt2[r]
            r += 1
            j += 1
        if(check_column(Histogram[i][0], 1) == 0):
            Histogram[i][0] = 'Select Column'
            messagebox.showerror("Incompatible Template","Template 3 is incompatible with the given data file. Please select the graphs manually")
            return
        i += 1
    
    # Violin (r = 3)
    i = 0
    while(i < 2):
        j = 0
        while(j < 3):
            if(j == 2):
                Violin[i][j] = int(listTxt2[r])
            else:
                Violin[i][j] = listTxt2[r]
            r += 1
            j += 1
        if(check_column(Violin[i][0], 1) == 0):
            Violin[i][0] = 'Select Column'
            messagebox.showerror("Incompatible Template","Template 3 is incompatible with the given data file. Please select the graphs manually")
            return
        i += 1
    # Pie (r = 9)
    i = 0
    while(i < 2):
        j = 0
        while(j < 4):
            if(j == 2 or j == 3):
                Pie[i][j] = int(listTxt2[r])
            else:
                Pie[i][j] = listTxt2[r]
            r += 1
            j += 1
        if(check_column(Pie[i][0], 2) == 0):
            Pie[i][0] = 'Select Column'
            messagebox.showerror("Incompatible Template","Template 3 is incompatible with the given data file. Please select the graphs manually")
            return
        i += 1
    # Count (r = 17)
    i = 0
    while(i < 2):
        j = 0
        while(j < 3):
            if(j == 2):
                Count[i][j] = int(listTxt2[r])
            else:
                Count[i][j] = listTxt2[r]
            r += 1
            j += 1
        if(check_column(Count[i][0], 2) == 0):
            Count[i][0] = 'Select Column'
            messagebox.showerror("Incompatible Template","Template 3 is incompatible with the given data file. Please select the graphs manually")
            return
        i += 1
    # Bar (r = 23)
    i = 0
    while(i < 2):
        j = 0
        while(j < 7):
            if(j == 3 or j == 4 or j == 6):
                Bar[i][j] = int(listTxt2[r])
            else:
                Bar[i][j] = listTxt2[r]
            r += 1
            j += 1
        if(check_column(Bar[i][0], 2) == 0 or check_column(Bar[i][1], 1) == 0 or check_column(Bar[i][5], 2) == 0):
            Bar[i][0] = 'Select Column'
            Bar[i][1] = 'Select Column'
            messagebox.showerror("Incompatible Template","Template 3 is incompatible with the given data file. Please select the graphs manually")
            return
        i += 1
    # Line (r = 37)
    i = 0
    while(i < 2):
        j = 0
        while(j < 7):
            if(j == 3 or j == 5 or j == 6):
                Line[i][j] = int(listTxt2[r])
            else:
                Line[i][j] = listTxt2[r]
            r += 1
            j += 1
        if(LinearSearch_YesNo(numerical_data + timeseries_data, Line[i][0]) == 0 or LinearSearch_YesNo(numerical_data + timeseries_data, Line[i][1]) == 0 or LinearSearch_YesNo(numerical_data + timeseries_data, Line[i][4]) == 0):
            ''''For Line Graph only: If datatype is converted from Categorical to others it cannot be stored in the template and manual entry needed'''
            print(Line)
            Line[i][0] = 'Select Column'
            Line[i][1] = 'Select Column'
            Line[i][4] = 'Select Column'
            messagebox.showerror("Incompatible Template","Template 3 is incompatible with the given data file. Please select the graphs manually")
            return
        i += 1
        
    messagebox.showinfo("Template Loaded","Template 3 has been loaded successfully.")
    return
        
def backend():
    '''This is the main function responsible for creating the presentation from the given data '''
    window_title.destroy()

    # Save Template if asked by user
    if(saveTemplate.get() != 0):
        createTemplate(saveTemplate.get())
        # Add template's name to AppFinal_TemplateNames.txt
        listTxt2 = []
        try:
            with open("AppFinal_TemplateNames.txt") as f:
                for l in f.readlines(): # Goes line by line
                    listTxt2.append(l.replace('\n',''))
        except FileNotFoundError:
            listTxt2 = ["Template 1","Template 2","Template 3"]
            listTxt2[saveTemplate.get() - 1] = templateName.get()
            with open("AppFinal_TemplateNames.txt", "w") as f:
                for l in listTxt2:
                    print(l, file = f)
        else:
            listTxt2[saveTemplate.get() - 1] = templateName.get()
            with open("AppFinal_TemplateNames.txt", "w") as f:
                for l in listTxt2:
                    print(l, file = f)
            
    
    pr = Presentation('default.pptx')
    
    # Title Slide
    if(presentationTitle.get() != "" or presentationSubtitle.get() != ""): # If Title or Subtitle exists
        slide_register_title = pr.slide_layouts[0]
        slide_title = pr.slides.add_slide(slide_register_title)
        title1 = slide_title.shapes.title
        title1.text = presentationTitle.get()

        subtitle1 = slide_title.placeholders[1]
        subtitle1.text = presentationSubtitle.get()
        PresFilename = presentationTitle.get() + ".pptx"
        if(PresFilename == "default.pptx"): # Can't have file name default as that name belongs to input of Presentation()
            PresFilename = "Presentation " + now.strftime("%d_%m_%Y %H_%M_%S") +".pptx"
    else:
        # No title slide if no title or subtitle is given
        now = datetime.now()
        PresFilename = "Presentation " + now.strftime("%d_%m_%Y %H_%M_%S") +".pptx" # Presentation named after timestamp of its creation

    
    # Histogram
    slide_register_hist = [pr.slide_layouts[5], pr.slide_layouts[5]]
    slide_hist = []
    title_hist = []
    
    i = 0
    e = 0
    while(i < 2):
        if(Histogram[i][0] == "Select Column"): # Redundancy for foolproofing
            i += 1
            continue
        slide_hist.append(pr.slides.add_slide(slide_register_hist[e]))
        title_hist.append(slide_hist[e].shapes.title)
        title_hist[e].text = Histogram[i][1]
        hist_Col_1 = pd.DataFrame(df[Histogram[i][0]])
        hist_Col_1 = hist_Col_1.dropna(how='all', axis = 0)
        f = plt.figure(figsize=(8, 4))
        sb.histplot(data = hist_Col_1)
        #plt.show()
        plt.savefig('graph1.jpg')
        img1 = 'graph1.jpg'
        add_picture = slide_hist[e].shapes.add_picture(img1, Inches(1), Inches(2))
        i += 1
        e += 1

    # Violin Plot
    slide_register_vio = [pr.slide_layouts[5], pr.slide_layouts[5]]
    slide_vio = []
    title_vio = []
    
    i = 0
    e = 0
    while(i < 2):
        if(Violin[i][0] == "Select Column"): # Redundancy for foolproofing
            i += 1
            continue
        slide_vio.append(pr.slides.add_slide(slide_register_vio[e]))
        title_vio.append(slide_vio[e].shapes.title)
        title_vio[e].text = Violin[i][1]
        vio_Col_1 = pd.DataFrame(df[Violin[i][0]])
        vio_Col_1 = vio_Col_1.dropna(how='all', axis = 0)
        f = plt.figure(figsize=(9, 4.5))
        sb.violinplot(data = vio_Col_1, orient = "h", color = colour_list[0][Violin[i][2]]) # Graph custom colour
        #plt.show()
        plt.savefig('graph1.jpg')
        img1 = 'graph1.jpg'
        add_picture = slide_vio[e].shapes.add_picture(img1, Inches(0.5), Inches(2))
        i += 1
        e += 1

    # Pie Plot
    slide_register_pie = [pr.slide_layouts[5], pr.slide_layouts[5]]
    slide_pie = []
    title_pie = []
    
    i = 0
    e = 0
    while(i < 2):
        if(Pie[i][0] == "Select Column"): # Redundancy for foolproofing
            i += 1
            continue
        slide_pie.append(pr.slides.add_slide(slide_register_pie[e]))
        title_pie.append(slide_pie[e].shapes.title)
        title_pie[e].text = Pie[i][1]

        # Making the Graph
        df_pie = pd.DataFrame(df[Pie[i][0]]).dropna(axis = 0)
        names = df_pie[Pie[i][0]].unique().tolist() # A list of all categories
        df_pie2 = df_pie[Pie[i][0]].value_counts().to_frame()
        # Finding the number of each category in names
        i1 = 0
        nos = []
        while(i1 < len(names)):
            nos.append(df_pie2[Pie[i][0]][names[i1]])
            i1 += 1

        # Bubble Sort - Descending Order for numbers
        i1 = 0
        j = 0
        while(i1 < len(names)):
            j = i1+1
            while(j < len(names)):
                if(nos[j] > nos[i1]):
                    temp = nos[j]
                    temp2 = names[j]
                    nos[j] = nos[i1]
                    names[j] = names[i1] # Names must be kept with their corresponding numbers
                    nos[i1] = temp
                    names[i1] = temp2
                j += 1
            i1 += 1

        # Optimising pie graph (Use of others if large of of slices such that 'Others' < 10%)
        max_slices = pie_find_max_slices(nos)
        if(Pie[i][2] == 1 and len(names) > max_slices): # Only if user wants and there are large number of slices
            names[max_slices - 1] = 'Others'
            i1 = max_slices
            while(i1 < len(nos)):
                nos[max_slices - 1] += nos[i1]
                i1 += 1
            i1 = len(nos) - 1
            while(i1 > max_slices - 1):
                nos.pop()
                names.pop()
                i1 -= 1
        
        # Label distance: gives the space between labels and the center of the pie
        f = plt.figure(figsize=(10, 5))
        plt.pie(nos, labels=names, labeldistance=1.15,autopct='%1.1f%%', colors = sb.color_palette(colour_list[2][Pie[i][3]]))
        #plt.show()

        plt.savefig('graph1.jpg')
        img1 = 'graph1.jpg'
        add_picture = slide_pie[e].shapes.add_picture(img1, Inches(0), Inches(2))
        i += 1
        e += 1

    # Count Plot
    slide_register_cou = [pr.slide_layouts[5], pr.slide_layouts[5]]
    slide_cou = []
    title_cou = []
    
    i = 0
    e = 0
    while(i < 2):
        if(Count[i][0] == "Select Column"): # Redundancy for foolproofing
            i += 1
            continue
        slide_cou.append(pr.slides.add_slide(slide_register_cou[e]))
        title_cou.append(slide_cou[e].shapes.title)
        title_cou[e].text = Count[i][1]
        cou_Col_1 = pd.DataFrame(df[Count[i][0]])
        cou_Col_1 = cou_Col_1.dropna(how='all', axis = 0)
        f = plt.figure(figsize=(10, 6))
        sb.catplot(y = Count[i][0], data = cou_Col_1, kind = "count", color = colour_list[0][Count[i][2]])
        #plt.show()
        plt.savefig('graph1.jpg')
        img1 = 'graph1.jpg'
        add_picture = slide_cou[e].shapes.add_picture(img1, Inches(0.5), Inches(2))
        i += 1
        e += 1

    # Bar Plot
    slide_register_bar = [pr.slide_layouts[5], pr.slide_layouts[5]]
    slide_bar = []
    title_bar = []
    
    i = 0
    e = 0
    while(i < 2):
        if(Bar[i][0] == "Select Column"): # Redundancy for foolproofing
            i += 1
            continue
        slide_bar.append(pr.slides.add_slide(slide_register_bar[e]))
        title_bar.append(slide_bar[e].shapes.title)
        title_bar[e].text = Bar[i][2]
        
        f = plt.figure(figsize=(10, 5))
        
        # plot a bar chart
        if(Bar[i][5] == 'Select Column'): # Not a Grouped Bar Chart
            if(Bar[i][3] == 0 and Bar[i][4] == 0):
                sb.barplot(x = df[Bar[i][0]], y = df[Bar[i][1]], data = df, estimator = sum, ci = None, color = colour_list[0][Bar[i][6]])
            elif(Bar[i][3] == 1 and Bar[i][4] == 0):
                sb.barplot(x = df[Bar[i][0]], y = df[Bar[i][1]], data = df, estimator = np.mean, ci = None, color = colour_list[0][Bar[i][6]])    
            elif(Bar[i][3] == 0 and Bar[i][4] == 1):
                sb.barplot(x = df[Bar[i][1]], y = df[Bar[i][0]], data = df, estimator = sum, ci = None, color = colour_list[0][Bar[i][6]])
            elif(Bar[i][3] == 1 and Bar[i][4] == 1):
                sb.barplot(x = df[Bar[i][1]], y = df[Bar[i][0]], data = df, estimator = np.mean, ci = None, color = colour_list[0][Bar[i][6]])
        else: # Grouped Bar Chart
            if(df[Bar[i][5]].value_counts().count() <= 2):
                barPlotColours = ['','']
                barPlotColours[0] = colour_list[0][Bar[i][6]]
                barPlotColours[1] = colour_list[1][Bar[i][6]]
                sb.set_palette(sb.color_palette(barPlotColours))
            else:
                sb.set_palette(sb.color_palette(colour_list[2][Bar[i][6]]))
            if(Bar[i][3] == 0 and Bar[i][4] == 0):
                sb.barplot(x = df[Bar[i][0]], y = df[Bar[i][1]], hue=df[Bar[i][5]], data = df, ci=None, estimator = sum)
            elif(Bar[i][3] == 1 and Bar[i][4] == 0):
                sb.barplot(x = df[Bar[i][0]], y = df[Bar[i][1]], hue=df[Bar[i][5]], data = df, ci=None, estimator = np.mean)    
            elif(Bar[i][3] == 0 and Bar[i][4] == 1):
                sb.barplot(x = df[Bar[i][1]], y = df[Bar[i][0]], hue=df[Bar[i][5]], data = df, ci=None, estimator = sum)
            elif(Bar[i][3] == 1 and Bar[i][4] == 1):
                sb.barplot(x = df[Bar[i][1]], y = df[Bar[i][0]], hue=df[Bar[i][5]], data = df, ci=None, estimator = np.mean)

        if(Bar[i][3] == 0):
            plt.ylabel(Bar[i][1] + " (Sum)")
        else:
            plt.ylabel(Bar[i][1] + " (Mean)")

        #plt.show()

        plt.savefig('graph1.jpg')
        img1 = 'graph1.jpg'
        add_picture = slide_bar[e].shapes.add_picture(img1, Inches(0), Inches(2))
        i += 1
        e += 1
    
    # Line Plot
    slide_register_lin = [pr.slide_layouts[5], pr.slide_layouts[5]]
    slide_lin = []
    title_lin = []
    
    i = 0
    e = 0
    while(i < 2):
        if(Line[i][0] == "Select Column"): # Redundancy for foolproofing
            i += 1
            continue
        slide_lin.append(pr.slides.add_slide(slide_register_lin[e]))
        title_lin.append(slide_lin[e].shapes.title)
        title_lin[e].text = Line[i][2]
        
        f = plt.figure(figsize=(10, 5))

        #Plotting a line graph
        if(Line[i][3] == 0):
            plt.plot(Line[i][0], Line[i][1], data=df, color = colour_list[0][Line[i][6]], linestyle = '-')
        elif(Line[i][3] == 1):
            plt.plot(Line[i][0], Line[i][1], data=df, color = colour_list[0][Line[i][6]], linestyle = '--')
        elif(Line[i][3] == 2):
            plt.plot(Line[i][0], Line[i][1], data=df, color = colour_list[0][Line[i][6]], linestyle = ':')
        else:
            plt.plot(Line[i][0], Line[i][1], data=df, color = colour_list[0][Line[i][6]], linestyle = '-.')
        #Plotting 2nd line graph
        if(Line[i][4] != 'Select Column'):
            if(Line[i][5] == 0):
                plt.plot(Line[i][0], Line[i][4], data=df, color = colour_list[1][Line[i][6]], linestyle = '-')
            elif(Line[i][5] == 1):
                plt.plot(Line[i][0], Line[i][4], data=df, color = colour_list[1][Line[i][6]], linestyle = '--')
            elif(Line[i][5] == 2):
                plt.plot(Line[i][0], Line[i][4], data=df, color = colour_list[1][Line[i][6]], linestyle = ':')
            else:
                plt.plot(Line[i][0], Line[i][4], data=df, color = colour_list[1][Line[i][6]], linestyle = '-.')

        plt.xlabel(Line[i][0])
        plt.legend()
        #plt.show()

        plt.savefig('graph1.jpg')
        img1 = 'graph1.jpg'
        add_picture = slide_lin[e].shapes.add_picture(img1, Inches(0), Inches(2))
        i += 1
        e += 1
        
    # Saving the presentation
    pr.save(PresFilename)
    os.remove("graph1.jpg")
    window.destroy()
    os.startfile(PresFilename)
    import sys
    sys.exit() # To break the mainloop in window which would still be running



# Space for handle click functions (frontend) starts

def title_slide(event):
    '''Function Triggered when Create Presentation is clicked and takes input of title slide details '''
    global window_title
    window_title = tk.Toplevel(window)
    window_title.title("Title Slide")

    global saveTemplate
    saveTemplate = tk.IntVar()
    saveTemplate.set(0) # 0 -> no save, 1 -> Save as template 1.

    global templateName
    templateName = tk.StringVar()
    templateName.set("Template 1")

    tk.Label(master = window_title, text = "Presentation Title").grid(row = 0, column = 0)
    tk.Label(master = window_title, text = "Subtitle (if any)").grid(row = 1, column = 0)
    tk.Entry(master = window_title, textvariable = presentationTitle, bg="white", width=20).grid(row = 0, column = 2)
    tk.Entry(master = window_title, textvariable = presentationSubtitle, bg="white", width=20).grid(row = 1, column = 2)
    tk.Radiobutton(window_title, text = "Dont Save Template", variable = saveTemplate, value = 0).grid(row = 2, column = 0)
    tk.Radiobutton(window_title, text = "Save as Template 1", variable = saveTemplate, value = 1).grid(row = 2, column = 2)
    tk.Radiobutton(window_title, text = "Save as Template 2", variable = saveTemplate, value = 2).grid(row = 3, column = 0)
    tk.Radiobutton(window_title, text = "Save as Template 3", variable = saveTemplate, value = 3).grid(row = 3, column = 2)
    tk.Label(master = window_title, text = "Template Name").grid(row = 4, column = 0)
    tk.Entry(master = window_title, textvariable = templateName, bg="white", width=20).grid(row = 4, column = 2)
    tk.Label(master = window_title, text = " ").grid(row = 0, column = 1)
    tk.Label(master = window_title, text = " ").grid(row = 5, column = 1)
    btn_createPres = tk.Button(master = window_title,  text = "Create Presentation", width=16,height=2,bg="white",fg="red")
    btn_createPres.bind("<Button-1>", check_all_inputs)
    btn_createPres.grid(row = 6, column = 2)
    

def check_all_inputs(event):
    '''Checks all inputs and passes control to backend if all is clear else throws an error.
       Triggered from window_title'''
    all_inputs_valid = 1 # i.e True
    #Histogram
    i = 0
    while(i < 2):
        if(LinearSearch_YesNo(numerical_data,Histogram[i][0]) == 0):
            all_inputs_valid = 0 #i.e false (invalid input exists)
            error_line = "Data Type Error:\n" + Histogram[i][0] + " is not a valid entry for Histogram " + str(i+1) + ". Please change its datatype to Number"
            messagebox.showwarning("Data Type Error",error_line)
            return
        i += 1
    #Violin
    i = 0
    while(i < 2):
        if(LinearSearch_YesNo(numerical_data,Violin[i][0]) == 0):
            all_inputs_valid = 0 #i.e false (invalid input exists)
            error_line = "Data Type Error:\n" + Violin[i][0] + " is not a valid entry for Violin Chart " + str(i+1) + ". Please change its datatype to Number"
            messagebox.showwarning("Data Type Error",error_line)

            return
        i += 1
    #Pie
    i = 0
    while(i < 2):
        if(LinearSearch_YesNo(categorical_data,Pie[i][0]) == 0):
            all_inputs_valid = 0 #i.e false (invalid input exists)
            error_line = "Data Type Error:\n" + Pie[i][0] + " is not a valid entry for Pie Chart " + str(i+1) + ". Please change its datatype to Category/Object"
            messagebox.showwarning("Data Type Error",error_line)
            return
        i += 1
    #Count
    i = 0
    while(i < 2):
        if(LinearSearch_YesNo(categorical_data,Count[i][0]) == 0):
            all_inputs_valid = 0 #i.e false (invalid input exists)
            error_line = "Data Type Error:\n" + Count[i][0] + " is not a valid entry for Count Chart " + str(i+1) + ". Please change its datatype to Category/Object"
            messagebox.showwarning("Data Type Error",error_line)
            return
        i += 1
    #Bar
    i = 0
    while(i < 2):
        if(LinearSearch_YesNo(categorical_data,Bar[i][0]) == 0):
            all_inputs_valid = 0 #i.e false (invalid input exists)
            error_line = "Data Type Error:\n" + Bar[i][0] + " is not a valid entry for Bar Chart " + str(i+1) + ". Please change its datatype to Category/Object"
            messagebox.showwarning("Data Type Error",error_line)
            return
        if(LinearSearch_YesNo(numerical_data,Bar[i][1]) == 0):
            all_inputs_valid = 0 #i.e false (invalid input exists)
            error_line = "Data Type Error:\n" + Bar[i][1] + " is not a valid entry for Bar Chart " + str(i+1) + ". Please change its datatype to Number"
            messagebox.showwarning("Data Type Error",error_line)
            return
        if(LinearSearch_YesNo(categorical_data,Bar[i][5]) == 0): # For Hue
            all_inputs_valid = 0 #i.e false (invalid input exists)
            error_line = "Data Type Error:\n" + Bar[i][5] + " is not a valid entry for Bar Chart " + str(i+1) + ". Please change its datatype to Category/Object"
            messagebox.showwarning("Data Type Error",error_line)
            return
        i += 1

    if(all_inputs_valid == 1):
        backend()

def insert_histogram(event):
    '''Takes input data for a histogram and stores it in variables'''
    global window_histogram
    window_histogram = tk.Toplevel(window)
    window_histogram.title("Histogram Entry")

    global histCol1
    histCol1 = tk.StringVar()
    histCol1.set(Histogram[0][0])
    global histTi1
    histTi1 = tk.StringVar()
    histTi1.set(Histogram[0][1])
    global histCol2
    histCol2 = tk.StringVar()
    histCol2.set(Histogram[1][0])
    global histTi2
    histTi2 = tk.StringVar()
    histTi2.set(Histogram[1][1])

    lbl_hist1 = tk.Label(master = window_histogram, text = "Histogram 1")
    lbl_hist1_col = tk.Label(master = window_histogram, text = "Column name")
    lbl_hist1_ti = tk.Label(master = window_histogram, text = "Graph Title")
    ent_hist1_ti = tk.Entry(master = window_histogram, textvariable = histTi1, bg="white", width=20)
    drp_hist1_col = tk.OptionMenu(window_histogram, histCol1, "Select Column", *numerical_data)
    lbl_hist1.grid(row = 0, column = 0)
    lbl_hist1_col.grid(row = 1, column = 0)
    drp_hist1_col.grid(row = 1, column = 1)
    lbl_hist1_ti.grid(row = 2, column = 0)
    ent_hist1_ti.grid(row = 2, column = 1)

    tk.Label(master = window_histogram, text = " ").grid(row = 3, column = 0)

    lbl_hist2 = tk.Label(master = window_histogram, text = "Histogram 2")
    lbl_hist2_col = tk.Label(master = window_histogram, text = "Column name")
    lbl_hist2_ti = tk.Label(master = window_histogram, text = "Graph Title")
    ent_hist2_ti = tk.Entry(master = window_histogram, textvariable = histTi2, bg="white", width=20)
    drp_hist2_col = tk.OptionMenu(window_histogram, histCol2, "Select Column", *numerical_data)
    lbl_hist2.grid(row = 4, column = 0)
    lbl_hist2_col.grid(row = 5, column = 0)
    drp_hist2_col.grid(row = 5, column = 1)
    lbl_hist2_ti.grid(row = 6, column = 0)
    ent_hist2_ti.grid(row = 6, column = 1)

    tk.Label(master = window_histogram, text = " ").grid(row = 7, column = 0)

    btn_hist_Sub = tk.Button(master = window_histogram, text = "Submit", width=6,height=2,bg="white",fg="red")
    btn_hist_Sub.bind("<Button-1>", submit_histogram)
    btn_hist_Sub.grid(row = 8, column = 1)

    window_histogram.mainloop()

def submit_histogram(event):
    '''Takes the input data for histogram stored as tk variables and loads them in the list Histogram'''
    print("Submitting Histogram Data")
    Histogram[0][0] = histCol1.get()
    Histogram[0][1] = histTi1.get()
    Histogram[1][0] = histCol2.get()
    Histogram[1][1] = histTi2.get()
    print(Histogram)
    window_histogram.destroy()

def insert_violin(event):
    global window_violin
    window_violin = tk.Toplevel(window)
    window_violin.title("Violin Plot Entry")

    global vioCol1
    vioCol1 = tk.StringVar()
    vioCol1.set(Violin[0][0])
    global vioTi1
    vioTi1 = tk.StringVar()
    vioTi1.set(Violin[0][1])
    global vioRang1 # To store colour
    vioRang1 = tk.StringVar()
    vioRang1.set(colour_list[0][Violin[0][2]])
    global vioCol2
    vioCol2 = tk.StringVar()
    vioCol2.set(Violin[1][0])
    global vioTi2
    vioTi2 = tk.StringVar()
    vioTi2.set(Violin[1][1])
    global vioRang2
    vioRang2 = tk.StringVar()
    vioRang2.set(colour_list[0][Violin[1][2]])

    lbl_vio1 = tk.Label(master = window_violin, text = "Violin Graph 1")
    lbl_vio1_col = tk.Label(master = window_violin, text = "Column name")
    lbl_vio1_ti = tk.Label(master = window_violin, text = "Graph Title")
    ent_vio1_ti = tk.Entry(master = window_violin, textvariable = vioTi1, bg="white", width=20)
    drp_vio1_col = tk.OptionMenu(window_violin, vioCol1, "Select Column", *numerical_data)
    lbl_vio1.grid(row = 0, column = 0)
    lbl_vio1_col.grid(row = 1, column = 0)
    drp_vio1_col.grid(row = 1, column = 1)
    lbl_vio1_ti.grid(row = 2, column = 0)
    ent_vio1_ti.grid(row = 2, column = 1)
    lbl_vio1_rang = tk.Label(master = window_violin, text = "Graph Colour").grid(row = 3, column = 0)
    drp_vio1_rang = tk.OptionMenu(window_violin, vioRang1, *colour_list[0]).grid(row = 3, column = 1)

    tk.Label(master = window_violin, text = " ").grid(row = 4, column = 0)

    lbl_vio2 = tk.Label(master = window_violin, text = "Violin Graph 2")
    lbl_vio2_col = tk.Label(master = window_violin, text = "Column name")
    lbl_vio2_ti = tk.Label(master = window_violin, text = "Graph Title")
    ent_vio2_ti = tk.Entry(master = window_violin, textvariable = vioTi2, bg="white", width=20)
    drp_vio2_col = tk.OptionMenu(window_violin, vioCol2, "Select Column", *numerical_data)
    lbl_vio2.grid(row = 5, column = 0)
    lbl_vio2_col.grid(row = 6, column = 0)
    drp_vio2_col.grid(row = 6, column = 1)
    lbl_vio2_ti.grid(row = 7, column = 0)
    ent_vio2_ti.grid(row = 7, column = 1)
    lbl_vio2_rang = tk.Label(master = window_violin, text = "Graph Colour").grid(row = 8, column = 0)
    drp_vio2_rang = tk.OptionMenu(window_violin, vioRang2, *colour_list[0]).grid(row = 8, column = 1)
    
    tk.Label(master = window_violin, text = " ").grid(row = 9, column = 0)

    btn_vio_Sub = tk.Button(master = window_violin, text = "Submit", width=6,height=2,bg="white",fg="red")
    btn_vio_Sub.bind("<Button-1>", submit_violin)
    btn_vio_Sub.grid(row = 10, column = 1)

    window_violin.mainloop()

def submit_violin(event):
    '''Takes the input data for violin stored as tk variables and loads them in the list Violin'''
    print("Submitting Violin plot Data")
    Violin[0][0] = vioCol1.get()
    Violin[0][1] = vioTi1.get()
    if(vioRang1.get() == "purple"): # As the colour purple is represented by m in seaborn
        Violin[0][2] = 7
    else:
        Violin[0][2] = LinearSearch_index(colour_list[0], vioRang1.get())
    Violin[1][0] = vioCol2.get()
    Violin[1][1] = vioTi2.get()
    if(vioRang2.get() == "purple"): # As the colour purple is represented by m in seaborn
        Violin[1][2] = 7
    else:
        Violin[1][2] = LinearSearch_index(colour_list[0], vioRang2.get())
    print(Violin)
    window_violin.destroy()

def insert_pie(event):
    global window_pie
    window_pie = tk.Toplevel(window)
    window_pie.title("Pie Plot Entry")

    global pieCol1
    pieCol1 = tk.StringVar()
    pieCol1.set(Pie[0][0])
    global pieTi1
    pieTi1 = tk.StringVar()
    pieTi1.set(Pie[0][1])
    global pieSmrt1
    pieSmrt1 = tk.IntVar()
    pieSmrt1.set(Pie[0][2])
    global pieRang1
    pieRang1 = tk.StringVar()
    pieRang1.set(colour_list[0][Pie[0][3]])
    global pieCol2
    pieCol2 = tk.StringVar()
    pieCol2.set(Pie[1][0])
    global pieTi2
    pieTi2 = tk.StringVar()
    pieTi2.set(Pie[1][1])
    global pieSmrt2
    pieSmrt2 = tk.IntVar()
    pieSmrt2.set(Pie[1][2])
    global pieRang2
    pieRang2 = tk.StringVar()
    pieRang2.set(colour_list[0][Pie[1][3]])

    lbl_pie1 = tk.Label(master = window_pie, text = "Pie Plot 1")
    lbl_pie1_col = tk.Label(master = window_pie, text = "Column name")
    lbl_pie1_ti = tk.Label(master = window_pie, text = "Graph Title")
    ent_pie1_ti = tk.Entry(master = window_pie, textvariable = pieTi1, bg="white", width=20)
    drp_pie1_col = tk.OptionMenu(window_pie, pieCol1, "Select Column", *categorical_data)
    chk_pie1_smrt = tk.Checkbutton(window_pie, text = "Use Others", variable = pieSmrt1, onvalue = 1, offvalue = 0)
    lbl_pie1.grid(row = 0, column = 0)
    lbl_pie1_col.grid(row = 1, column = 0)
    drp_pie1_col.grid(row = 1, column = 1)
    lbl_pie1_ti.grid(row = 2, column = 0)
    ent_pie1_ti.grid(row = 2, column = 1)
    chk_pie1_smrt.grid(row = 3, column = 0)
    lbl_pie1_rang = tk.Label(master = window_pie, text = "Graph Colour").grid(row = 4, column = 0)
    drp_pie1_rang = tk.OptionMenu(window_pie, pieRang1, *colour_list[0]).grid(row = 4, column = 1)


    tk.Label(master = window_pie, text = " ").grid(row = 5, column = 0)

    lbl_pie2 = tk.Label(master = window_pie, text = "Pie Plot 2")
    lbl_pie2_col = tk.Label(master = window_pie, text = "Column name")
    lbl_pie2_ti = tk.Label(master = window_pie, text = "Graph Title")
    ent_pie2_ti = tk.Entry(master = window_pie, textvariable = pieTi2, bg="white", width=20)
    drp_pie2_col = tk.OptionMenu(window_pie, pieCol2, "Select Column", *categorical_data)
    chk_pie2_smrt = tk.Checkbutton(window_pie, text = "Use Others", variable = pieSmrt2, onvalue = 1, offvalue = 0)
    lbl_pie2.grid(row = 6, column = 0)
    lbl_pie2_col.grid(row = 7, column = 0)
    drp_pie2_col.grid(row = 7, column = 1)
    lbl_pie2_ti.grid(row = 8, column = 0)
    ent_pie2_ti.grid(row = 8, column = 1)
    chk_pie2_smrt.grid(row = 9, column = 0)
    lbl_pie2_rang = tk.Label(master = window_pie, text = "Graph Colour").grid(row = 10, column = 0)
    drp_pie2_rang = tk.OptionMenu(window_pie, pieRang2, *colour_list[0]).grid(row = 10, column = 1)

    tk.Label(master = window_pie, text = " ").grid(row = 11, column = 0)

    btn_pie_Sub = tk.Button(master = window_pie, text = "Submit", width=6,height=2,bg="white",fg="red")
    btn_pie_Sub.bind("<Button-1>", submit_pie)
    btn_pie_Sub.grid(row = 12, column = 1)

    window_pie.mainloop()

def submit_pie(event):
    '''Takes the input data for insert_pie stored as tk variables and loads them in the list Pie'''
    print("Submitting Pie plot Data")
    Pie[0][0] = pieCol1.get()
    Pie[0][1] = pieTi1.get()
    Pie[0][2] = pieSmrt1.get()
    if(pieRang1.get() == "purple"): # As the colour purple is represented by m in seaborn
        Pie[0][3] = 7
    else:
        Pie[0][3] = LinearSearch_index(colour_list[0], pieRang1.get())
    Pie[1][0] = pieCol2.get()
    Pie[1][1] = pieTi2.get()
    Pie[1][2] = pieSmrt2.get()
    if(pieRang2.get() == "purple"): # As the colour purple is represented by m in seaborn
        Pie[1][3] = 7
    else:
        Pie[1][3] = LinearSearch_index(colour_list[0], pieRang2.get())
    print(Pie)
    window_pie.destroy()

def insert_count(event):
    global window_count
    window_count = tk.Toplevel(window)
    window_count.title("Count Plot Entry")

    global couCol1
    couCol1 = tk.StringVar()
    couCol1.set(Count[0][0])
    global couTi1
    couTi1 = tk.StringVar()
    couTi1.set(Count[0][1])
    global couRang1
    couRang1 = tk.StringVar()
    couRang1.set(colour_list[0][Count[0][2]])
    global couCol2
    couCol2 = tk.StringVar()
    couCol2.set(Count[1][0])
    global couTi2
    couTi2 = tk.StringVar()
    couTi2.set(Count[1][1])
    global couRang2
    couRang2 = tk.StringVar()
    couRang2.set(colour_list[0][Count[1][2]])

    lbl_cou1 = tk.Label(master = window_count, text = "Count Plot 1")
    lbl_cou1_col = tk.Label(master = window_count, text = "Column name")
    lbl_cou1_ti = tk.Label(master = window_count, text = "Graph Title")
    ent_cou1_ti = tk.Entry(master = window_count, textvariable = couTi1, bg="white", width=20)
    drp_cou1_col = tk.OptionMenu(window_count, couCol1, "Select Column", *categorical_data)
    lbl_cou1.grid(row = 0, column = 0)
    lbl_cou1_col.grid(row = 1, column = 0)
    drp_cou1_col.grid(row = 1, column = 1)
    lbl_cou1_ti.grid(row = 2, column = 0)
    ent_cou1_ti.grid(row = 2, column = 1)
    tk.Label(master = window_count, text = "Graph Colour").grid(row = 3, column = 0)
    tk.OptionMenu(window_count, couRang1, *colour_list[0]).grid(row = 3, column = 1)

    tk.Label(master = window_count, text = " ").grid(row = 4, column = 0)

    lbl_cou2 = tk.Label(master = window_count, text = "Count Plot 2")
    lbl_cou2_col = tk.Label(master = window_count, text = "Column name")
    lbl_cou2_ti = tk.Label(master = window_count, text = "Graph Title")
    ent_cou2_ti = tk.Entry(master = window_count, textvariable = couTi2, bg="white", width=20)
    drp_cou2_col = tk.OptionMenu(window_count, couCol2, "Select Column", *categorical_data)
    lbl_cou2.grid(row = 5, column = 0)
    lbl_cou2_col.grid(row = 6, column = 0)
    drp_cou2_col.grid(row = 6, column = 1)
    lbl_cou2_ti.grid(row = 7, column = 0)
    ent_cou2_ti.grid(row = 7, column = 1)
    tk.Label(master = window_count, text = "Graph Colour").grid(row = 8, column = 0)
    tk.OptionMenu(window_count, couRang2, *colour_list[0]).grid(row = 8, column = 1)
    
    tk.Label(master = window_count, text = " ").grid(row = 9, column = 0)

    btn_cou_Sub = tk.Button(master = window_count, text = "Submit", width=6,height=2,bg="white",fg="red")
    btn_cou_Sub.bind("<Button-1>", submit_count)
    btn_cou_Sub.grid(row = 10, column = 1)

    window_count.mainloop()

def submit_count(event):
    '''Takes the input data for count chart stored as tk variables and loads them in the list Count'''
    print("Submitting Count plot Data")
    Count[0][0] = couCol1.get()
    Count[0][1] = couTi1.get()
    if(couRang1.get() == "purple"): # As the colour purple is represented by m in seaborn
        Count[0][2] = 7
    else:
        Count[0][2] = LinearSearch_index(colour_list[0], couRang1.get())
    Count[1][0] = couCol2.get()
    Count[1][1] = couTi2.get()
    if(couRang2.get() == "purple"): # As the colour purple is represented by m in seaborn
        Count[1][2] = 7
    else:
        Count[1][2] = LinearSearch_index(colour_list[0], couRang2.get())
    print(Count)
    window_count.destroy()

def insert_bar(event):
    global window_bar
    window_bar = tk.Toplevel(window)
    window_bar.title("Bar Plot Entry")

    global barColCat1
    barColCat1 = tk.StringVar()
    barColCat1.set(Bar[0][0])
    global barColNum1
    barColNum1 = tk.StringVar()
    barColNum1.set(Bar[0][1])
    global barTi1
    barTi1 = tk.StringVar()
    barTi1.set(Bar[0][2])
    global barType1
    barType1 = tk.IntVar()
    barType1.set(Bar[0][3])
    global barInvert1
    barInvert1 = tk.IntVar()
    barInvert1.set(Bar[0][4])
    global barHue1
    barHue1 = tk.StringVar()
    barHue1.set(Bar[0][5])
    global barRang1
    barRang1 = tk.StringVar()
    barRang1.set(colour_list[0][Bar[0][6]])

    global barColCat2
    barColCat2 = tk.StringVar()
    barColCat2.set(Bar[1][0])
    global barColNum2
    barColNum2 = tk.StringVar()
    barColNum2.set(Bar[1][1])
    global barTi2
    barTi2 = tk.StringVar()
    barTi2.set(Bar[1][2])
    global barType2
    barType2 = tk.IntVar()
    barType2.set(Bar[1][3])
    global barInvert2
    barInvert2 = tk.IntVar()
    barInvert2.set(Bar[1][4])
    global barHue2
    barHue2 = tk.StringVar()
    barHue2.set(Bar[1][5])
    global barRang2
    barRang2 = tk.StringVar()
    barRang2.set(colour_list[0][Bar[1][6]])

    tk.Label(master = window_bar, text = "Bar Plot 1",font = "bold").grid(row = 0, column = 0)
    tk.Label(master = window_bar, text = "Column name (X-Axis)").grid(row = 1, column = 0)
    tk.OptionMenu(window_bar, barColCat1, "Select Column", *categorical_data).grid(row = 1, column = 1)
    tk.Label(master = window_bar, text = "Column name (Y-Axis)").grid(row = 2, column = 0)
    tk.OptionMenu(window_bar, barColNum1, "Select Column", *numerical_data).grid(row = 2, column = 1)
    tk.Radiobutton(window_bar, text = "Sum", variable = barType1, value = 0).grid(row = 3, column = 0)
    tk.Radiobutton(window_bar, text = "Mean", variable = barType1, value = 1).grid(row = 3, column = 1)
    chk_bar1_invert = tk.Checkbutton(window_bar, text = "Invert X & Y Axis", variable = barInvert1, onvalue = 1, offvalue = 0)
    chk_bar1_invert.grid(row = 4, column = 0)
    tk.Label(master = window_bar, text = "Graph Title").grid(row = 5, column = 0)
    tk.Entry(master = window_bar, textvariable = barTi1, bg="white", width=20).grid(row = 5, column = 1)
    tk.Label(master = window_bar, text = "Hue (if any)").grid(row =6,column = 0)
    tk.OptionMenu(window_bar, barHue1, "Select Column", *categorical_data).grid(row = 6, column = 1)
    tk.Label(master = window_bar, text = "Graph Colour").grid(row = 7, column = 0)
    tk.OptionMenu(window_bar, barRang1, *colour_list[0]).grid(row = 7, column = 1)

    tk.Label(master = window_bar, text = "").grid(row = 8, column = 0)

    tk.Label(master = window_bar, text = "Bar Plot 2",font = "bold").grid(row = 9, column = 0)
    tk.Label(master = window_bar, text = "Column name (X-Axis)").grid(row = 10, column = 0)
    tk.OptionMenu(window_bar, barColCat2, "Select Column", *categorical_data).grid(row = 10, column = 1)
    tk.Label(master = window_bar, text = "Column name (Y-Axis)").grid(row = 11, column = 0)
    tk.OptionMenu(window_bar, barColNum2, "Select Column", *numerical_data).grid(row = 11, column = 1)
    tk.Radiobutton(window_bar, text = "Sum", variable = barType2, value = 0).grid(row = 12, column = 0)
    tk.Radiobutton(window_bar, text = "Mean", variable = barType2, value = 1).grid(row = 12, column = 1)
    chk_bar2_invert = tk.Checkbutton(window_bar, text = "Invert X & Y Axis", variable = barInvert2, onvalue = 1, offvalue = 0)
    chk_bar2_invert.grid(row = 13, column = 0)
    tk.Label(master = window_bar, text = "Graph Title").grid(row = 14, column = 0)
    tk.Entry(master = window_bar, textvariable = barTi2, bg="white", width=20).grid(row = 14, column = 1)
    tk.Label(master = window_bar, text = "Hue (if any)").grid(row =15,column = 0)
    tk.OptionMenu(window_bar, barHue2, "Select Column", *categorical_data).grid(row = 15, column = 1)
    tk.Label(master = window_bar, text = "Graph Colour").grid(row = 16, column = 0)
    tk.OptionMenu(window_bar, barRang2, *colour_list[0]).grid(row = 16, column = 1)

    tk.Label(master = window_bar, text = "").grid(row = 17, column = 0)

    btn_bar_sub = tk.Button(master = window_bar, text = "Submit", width=6,height=2,bg="white",fg="red")
    btn_bar_sub.bind("<Button-1>", submit_bar)
    btn_bar_sub.grid(row = 18, column = 1)

    window_bar.mainloop()
    
def submit_bar(event):
    '''Takes the input data for Bar Chart stored as tk variables and loads them in the list Bar also raises errors if needed'''
    print("Submitting Bar Graph")
    if(barColCat1.get() == "Select Column" and barColNum1.get() != "Select Column"):
        messagebox.showerror("Data Missing","Please select the X axis data for Bar Plot 1.\nOtherwise, please use a Histogram.")
    elif(barColNum1.get() == "Select Column" and barColCat1.get() != "Select Column"):
        messagebox.showerror("Data Missing","Please select the Y axis data for Bar Plot 1.\nOtherwise, please use a Count Plot.")
    elif(barColCat2.get() == "Select Column" and barColNum2.get() != "Select Column"):
        messagebox.showerror("Data Missing","Please select the X axis data for Bar Plot 2.\nOtherwise, please use a Histogram.")
    elif(barColNum2.get() == "Select Column" and barColCat2.get() != "Select Column"):
        messagebox.showerror("Data Missing","Please select the Y axis data for Bar Plot 2.\nOtherwise, please use a Count Plot.")
    else:
        Bar[0][0] = barColCat1.get()
        Bar[0][1] = barColNum1.get()
        Bar[0][2] = barTi1.get()
        Bar[0][3] = barType1.get()
        Bar[0][4] = barInvert1.get()
        Bar[0][5] = barHue1.get()
        if(barRang1.get() == "purple"): # As the colour purple is represented by m in seaborn
            Bar[0][6] = 7
        else:
            Bar[0][6] = LinearSearch_index(colour_list[0], barRang1.get())
        Bar[1][0] = barColCat2.get()
        Bar[1][1] = barColNum2.get()
        Bar[1][2] = barTi2.get()
        Bar[1][3] = barType2.get()
        Bar[1][4] = barInvert2.get()
        Bar[1][5] = barHue2.get()
        if(barRang2.get() == "purple"): # As the colour purple is represented by m in seaborn
            Bar[1][6] = 7
        else:
            Bar[1][6] = LinearSearch_index(colour_list[0], barRang2.get())
        window_bar.destroy()
    print(Bar)

def insert_line(event):
    global window_lin
    window_lin = tk.Toplevel(window)
    window_lin.title("Line Plot Entry")

    global linColX1
    linColX1 = tk.StringVar()
    linColX1.set(Line[0][0])
    global linColY1
    linColY1 = tk.StringVar()
    linColY1.set(Line[0][1])
    global linTi1
    linTi1 = tk.StringVar()
    linTi1.set(Line[0][2])
    global linType1
    linType1 = tk.IntVar()
    linType1.set(Line[0][3])
    global linColY1_2
    linColY1_2 = tk.StringVar()
    linColY1_2.set(Line[0][4])
    global linType1_2
    linType1_2 = tk.IntVar()
    linType1_2.set(Line[0][5])
    global linRang1
    linRang1 = tk.StringVar()
    linRang1.set(colour_list[0][Line[0][6]])

    global linColX2
    linColX2 = tk.StringVar()
    linColX2.set(Line[1][0])
    global linColY2
    linColY2 = tk.StringVar()
    linColY2.set(Line[1][1])
    global linTi2
    linTi2 = tk.StringVar()
    linTi2.set(Line[1][2])
    global linType2
    linType2 = tk.IntVar()
    linType2.set(Line[1][3])
    global linColY2_2
    linColY2_2 = tk.StringVar()
    linColY2_2.set(Line[1][4])
    global linType2_2
    linType2_2 = tk.IntVar()
    linType2_2.set(Line[1][5])
    global linRang2
    linRang2 = tk.StringVar()
    linRang2.set(colour_list[0][Line[1][6]])

    tk.Label(master = window_lin, text = "Line Plot 1",font = "bold").grid(row = 0, column = 0)
    tk.Label(master = window_lin, text = "Column name (X-Axis)").grid(row = 1, column = 0)
    tk.OptionMenu(window_lin, linColX1, "Select Column", *numerical_data, *timeseries_data).grid(row = 1, column = 1)
    tk.Label(master = window_lin, text = "Column name (Y-Axis)").grid(row = 2, column = 0)
    tk.OptionMenu(window_lin, linColY1, "Select Column", *numerical_data, *timeseries_data).grid(row = 2, column = 1)
    tk.Radiobutton(window_lin, text = "Solid Line", variable = linType1, value = 0).grid(row = 3, column = 0)
    tk.Radiobutton(window_lin, text = "Dashed Line", variable = linType1, value = 1).grid(row = 3, column = 1)
    tk.Radiobutton(window_lin, text = "Dotted Line", variable = linType1, value = 2).grid(row = 4, column = 0)
    tk.Radiobutton(window_lin, text = "Dashed & Dotted Line", variable = linType1, value = 3).grid(row = 4, column = 1)
    tk.Label(master = window_lin, text = "Graph Title").grid(row = 5, column = 0)
    tk.Entry(master = window_lin, textvariable = linTi1, bg="white", width=20).grid(row = 5, column = 1)
    tk.Label(master = window_lin, text = "Column name 2 (Y-Axis)").grid(row = 6, column = 0)
    tk.OptionMenu(window_lin, linColY1_2, "Select Column", *numerical_data, *timeseries_data).grid(row = 6, column = 1)
    tk.Radiobutton(window_lin, text = "Solid Line", variable = linType1_2, value = 0).grid(row = 7, column = 0)
    tk.Radiobutton(window_lin, text = "Dashed Line", variable = linType1_2, value = 1).grid(row = 7, column = 1)
    tk.Radiobutton(window_lin, text = "Dotted Line", variable = linType1_2, value = 2).grid(row = 8, column = 0)
    tk.Radiobutton(window_lin, text = "Dashed & Dotted Line", variable = linType1_2, value = 3).grid(row = 8, column = 1)
    tk.Label(master = window_lin, text = "Graph Colour").grid(row = 9, column = 0)
    tk.OptionMenu(window_lin, linRang1, *colour_list[0]).grid(row = 9, column = 1)

    tk.Label(master = window_lin, text = "").grid(row = 10, column = 0)

    tk.Label(master = window_lin, text = "Line Plot 2",font = "bold").grid(row = 11, column = 0)
    tk.Label(master = window_lin, text = "Column name (X-Axis)").grid(row = 12, column = 0)
    tk.OptionMenu(window_lin, linColX2, "Select Column", *numerical_data, *timeseries_data).grid(row = 12, column = 1)
    tk.Label(master = window_lin, text = "Column name (Y-Axis)").grid(row = 13, column = 0)
    tk.OptionMenu(window_lin, linColY2, "Select Column", *numerical_data, *timeseries_data).grid(row = 13, column = 1)
    tk.Radiobutton(window_lin, text = "Solid Line", variable = linType2, value = 0).grid(row = 14, column = 0)
    tk.Radiobutton(window_lin, text = "Dashed Line", variable = linType2, value = 1).grid(row = 14, column = 1)
    tk.Radiobutton(window_lin, text = "Dotted Line", variable = linType2, value = 2).grid(row = 15, column = 0)
    tk.Radiobutton(window_lin, text = "Dashed & Dotted Line", variable = linType2, value = 3).grid(row = 15, column = 1)
    tk.Label(master = window_lin, text = "Graph Title").grid(row = 16, column = 0)
    tk.Entry(master = window_lin, textvariable = linTi2, bg="white", width=20).grid(row = 16, column = 1)
    tk.Label(master = window_lin, text = "Column name (Y-Axis)").grid(row = 17, column = 0)
    tk.OptionMenu(window_lin, linColY2_2, "Select Column", *numerical_data, *timeseries_data).grid(row = 17, column = 1)
    tk.Radiobutton(window_lin, text = "Solid Line", variable = linType2_2, value = 0).grid(row = 18, column = 0)
    tk.Radiobutton(window_lin, text = "Dashed Line", variable = linType2_2, value = 1).grid(row = 18, column = 1)
    tk.Radiobutton(window_lin, text = "Dotted Line", variable = linType2_2, value = 2).grid(row = 19, column = 0)
    tk.Radiobutton(window_lin, text = "Dashed & Dotted Line", variable = linType2_2, value = 3).grid(row = 19, column = 1)    
    tk.Label(master = window_lin, text = "Graph Colour").grid(row = 20, column = 0)
    tk.OptionMenu(window_lin, linRang2, *colour_list[0]).grid(row = 20, column = 1)
    
    tk.Label(master = window_lin, text = "").grid(row = 21, column = 0)

    btn_lin_sub = tk.Button(master = window_lin, text = "Submit", width=6,height=2,bg="white",fg="red")
    btn_lin_sub.bind("<Button-1>", submit_lin)
    btn_lin_sub.grid(row = 22, column = 1)

    window_lin.mainloop()

def submit_lin(event):
    '''Takes the input data for Line Chart stored as tk variables and loads them in the list Line, also raises errors if needed'''
    print("Submitting Line Graph")

    # For Line Chart 1
    if(linColX1.get() == "Select Column" and linColY1.get() == "Select Column"):
        i = 1 # Basically do nothing
    elif(linColX1.get() == "Select Column" and linColY1.get() != "Select Column"):
        messagebox.showerror("Data Missing","Please select the X axis data for Line Plot 1.\nOtherwise, please use a Histogram.")
    elif(linColY1.get() == "Select Column" and linColX1.get() != "Select Column"):
        messagebox.showerror("Data Missing","Please select the Y axis data for Bar Plot 1.\nOtherwise, please use a Histogram.")
    elif(LinearSearch_YesNo(numerical_data, linColY1.get()) == 1 and LinearSearch_YesNo(timeseries_data, linColY1_2.get()) == 1):
        messagebox.showwarning("Data Incompatible","Numerical & Time Series Data cannot be represented on the same axis. Please change the Y axis data of Line Plot 1.")
    elif(LinearSearch_YesNo(timeseries_data, linColY1.get()) == 1 and LinearSearch_YesNo(numerical_data, linColY1_2.get()) == 1):
        messagebox.showwarning("Data Incompatible","Numerical & Time Series Data cannot be represented on the same axis. Please change the Y axis data of Line Plot 1.")
    else:
        Line[0][0] = linColX1.get()
        Line[0][1] = linColY1.get()
        Line[0][2] = linTi1.get()
        Line[0][3] = linType1.get()
        Line[0][4] = linColY1_2.get()
        Line[0][5] = linType1_2.get()
        if(linRang1.get() == "purple"): # As the colour purple is represented by m in seaborn
            Line[0][6] = 7
        else:
            Line[0][6] = LinearSearch_index(colour_list[0], linRang1.get())

    # For Line Chart 2
    if(linColX2.get() == "Select Column" and linColY2.get() == "Select Column"):
        window_lin.destroy() # Basically do nothing about Line Chart 2
    elif(linColX2.get() == "Select Column" and linColY2.get() != "Select Column"):
        messagebox.showerror("Data Missing","Please select the X axis data for Bar Plot 2.\nOtherwise, please use a Histogram.")
    elif(linColY2.get() == "Select Column" and linColX2.get() != "Select Column"):
        messagebox.showerror("Data Missing","Please select the Y axis data for Bar Plot 2.\nOtherwise, please use a Histogram.")
    elif(LinearSearch_YesNo(numerical_data, linColY2.get()) == 1 and LinearSearch_YesNo(timeseries_data, linColY2_2.get()) == 1):
        messagebox.showwarning("Data Incompatible","Numerical & Time Series Data cannot be represented on the same axis. Please change the Y axis data of Line Plot 2.")
    
    elif(LinearSearch_YesNo(timeseries_data, linColY2.get()) == 1 and LinearSearch_YesNo(numerical_data, linColY2_2.get()) == 1):
        messagebox.showwarning("Data Incompatible","Numerical & Time Series Data cannot be represented on the same axis. Please change the Y axis data of Line Plot 2.")
    else:
        Line[1][0] = linColX2.get()
        Line[1][1] = linColY2.get()
        Line[1][2] = linTi2.get()
        Line[1][3] = linType2.get()
        Line[1][4] = linColY2_2.get()
        Line[1][5] = linType2_2.get()
        if(linRang1.get() == "purple"): # As the colour purple is represented by m in seaborn
            Line[1][6] = 7
        else:
            Line[1][6] = LinearSearch_index(colour_list[0], linRang2.get())
        window_lin.destroy()

    print(Line)

def Edit_Datatypes(event): #Front End for editing datatypes (takes input)
    global window_edit
    window_edit = tk.Toplevel(window)
    window_edit.title("Edit Data Types")

    global edit_column
    edit_column = tk.StringVar()
    edit_column.set("Select Column")
    global edit_newDataType
    edit_newDataType = tk.StringVar()
    edit_newDataType.set("Select Data Type")

    tk.Label(master = window_edit, text = " ").grid(row = 0, column = 0)
    tk.Label(master = window_edit, text = "Column Name").grid(row = 1, column = 0)
    tk.OptionMenu(window_edit, edit_column, "Select Column", *numerical_data, *timeseries_data, *categorical_data).grid(row = 1, column = 1)
    tk.Label(master = window_edit, text = "New Data Type").grid(row = 2, column = 0)
    tk.OptionMenu(window_edit, edit_newDataType, "Select Data Type", "Number", "Date & Time", "Object/Category").grid(row = 2, column = 1)
    tk.Label(master = window_edit, text = " ").grid(row = 3, column = 0)

    btn_edit_check = tk.Button(master = window_edit, text = "Check Data Type", width=12,height=1,bg="white",fg="blue")
    btn_edit_check.bind("<Button-1>", check_datatype)
    btn_edit_check.grid(row = 1, column = 2)

    btn_edit_con = tk.Button(master = window_edit, text = "Convert", width=6,height=2,bg="white",fg="red")
    btn_edit_con.bind("<Button-1>", convert_datatype)
    btn_edit_con.grid(row = 4, column = 1)

    window_edit.mainloop()

def check_datatype(event):
    if(edit_column.get() == 'Select Column'):
        messagebox.showerror("Error Data Absent","Error - Data Absent :-\nPlease select a column before checking its Data Type.")
    elif(LinearSearch_YesNo(numerical_data, edit_column.get()) == 1):
        messagebox.showinfo("Check Data Type","Column is a Number")
    elif(LinearSearch_YesNo(categorical_data, edit_column.get()) == 1):
        messagebox.showinfo("Check Data Type","Column is an Object/Category")
    elif(LinearSearch_YesNo(timeseries_data, edit_column.get()) == 1):
        messagebox.showinfo("Check Data Type","Column is a Date & Time")
    else:
        messagebox.showerror("Error Data Absent","Error - Data Absent :-\nPlease select a column before checking its Data Type.")

    window_edit.destroy()

def convert_datatype(event): #Back End for editing datatypes (actually edits the datatypes and raises eroors if needed)
    if(edit_column.get() == 'Select Column'):
        messagebox.showerror("Error Data Absent","Error - Data Absent :-\nPlease select a column before converting its Data Type.")
    elif(edit_newDataType.get() == 'Number'):
        # Not converting same datatype right?
        if(LinearSearch_YesNo(numerical_data, edit_column.get()) == 1):
            messagebox.showwarning("Convertion Failed","The Column is already of Number data type")
        else:
            # Now convert to numerical
            try:
                df[edit_column.get()] = pd.to_numeric(df[edit_column.get()])
            except ValueError as err:
                # if some error ask user if he wants to forcibly convert
                forcibly = messagebox.askyesno("Unable to convert","Unable to convert Column to Number. \n Would you like to do so forcibly?")
                if(forcibly == 1):
                    forcibly2 = messagebox.askokcancel("Warning","WARNING: Forcible conversion may result in loss of data. Do you still want to continue?")
                    if(forcibly2 == 1):
                        df[edit_column.get()] = pd.to_numeric(df[edit_column.get()], errors='coerce')
                        # Removing the column from its old type list
                        if(LinearSearch_YesNo(categorical_data, edit_column.get()) == 1):
                            categorical_data.remove(edit_column.get())
                        else:
                            timeseries_data.remove(edit_column.get())
                        # Add the same column to new type list
                        numerical_data.append(edit_column.get()) 
                        messagebox.showinfo("Conversion Succesful","Success: The column is now a Number") #Inform the user
                    else:
                        df[edit_column.get()] = pd.to_numeric(df[edit_column.get()], errors='ignore')
                else:
                    df[edit_column.get()] = pd.to_numeric(df[edit_column.get()], errors='ignore')
            else:
                '''If try block is executed'''
                # Removing the column from its old type list
                if(LinearSearch_YesNo(categorical_data, edit_column.get()) == 1):
                    categorical_data.remove(edit_column.get())
                else:
                    timeseries_data.remove(edit_column.get())
                # Add the same column to new type list
                numerical_data.append(edit_column.get()) 
                messagebox.showinfo("Conversion Succesful","Success: The column is now a Number") #Inform the user
    elif(edit_newDataType.get() == 'Date & Time'):
        if(LinearSearch_YesNo(timeseries_data, edit_column.get()) == 1):
            messagebox.showwarning("Conversion Failed","The Column is already of Date & Time data type")
        else:
            try:
                df[edit_column.get()] = df[edit_column.get()].astype('datetime64[ns]')
            except ValueError as err:
                messagebox.showinfo("Conversion Error","Failure: The column cannot be converted into Date & Time") #Inform the user
            else:
                '''If try block is executed'''
                # Removing the column from its old type list
                if(LinearSearch_YesNo(categorical_data, edit_column.get()) == 1):
                    categorical_data.remove(edit_column.get())
                else:
                    numerical_data.remove(edit_column.get())
                # Add the same column to new type list
                timeseries_data.append(edit_column.get()) 
                messagebox.showinfo("Conversion Succesful","Success: The column is now a Date & Time") #Inform the user
    elif(edit_newDataType.get() == 'Object/Category'):
        if(LinearSearch_YesNo(categorical_data, edit_column.get()) == 1):
            messagebox.showwarning("Convertion Failed","The Column is already of Date & Time data type")
        else:
            df[edit_column.get()] = df[edit_column.get()].astype('object')
            # Removing the column from its old type list
            if(LinearSearch_YesNo(timeseries_data, edit_column.get()) == 1):
                timeseries_data.remove(edit_column.get())
            else:
                numerical_data.remove(edit_column.get())
            # Add the same column to new type list
            categorical_data.append(edit_column.get()) 
            messagebox.showinfo("Conversion Succesful","Success: The column is now an Object/Category") #Inform the user
    else:
        messagebox.showerror("Error Data Absent","Error - Data Absent :-\nPlease select a New Data Type before converting")
                                
    window_edit.destroy()
    
def user_instructions(event):
    ''''Just a window that gives users instructions on how to use the application'''
    global window_instructions
    window_instructions = tk.Toplevel(window)
    window_instructions.title("X2P - Instructions")
    tk.Label(window_instructions,text = "Instructions", font = "bold").pack()
    tk.Label(window_instructions,text = "1) Click on the graph you want to insert ").pack()
    tk.Label(window_instructions,text = "2) If you cant find the column you want,  ").pack()
    tk.Label(window_instructions,text = "please check & change its datatype(if needed)").pack()
    tk.Label(window_instructions,text = "3) Click on the Templates (blue) to insert them").pack()
    tk.Label(window_instructions,text = "4) Click on Create Presentation to finish ").pack()
# Space for handle click functions (frontend) ends


window.filename = filedialog.askopenfilename(initialdir= "downloads", title = "Select Data File", filetypes=(("csv files","*.csv"),("xlsx files","*.xlsx")))
dataFile = window.filename

# Space for initial data analysis starts
if(dataFile[-4:] == ".csv"): # We have a csv file
    df = pd.read_csv(dataFile)
elif(dataFile[-5:] == ".xlsx"): 
    df = pd.read_excel(dataFile, sheet_name = 'Sheet1')
else:
    print("Error File type not supported")

df = df.dropna(how='all', axis = 1) # Removes all columns that are full of Null Values
columns = list(df) # A list of all columns in the dataframe df

numerical_data = [] # To store all columns with numerical data
categorical_data = []
timeseries_data = []
for i in columns:
    if(df.dtypes[i] == np.object):
        categorical_data.append(i)
    elif(df.dtypes[i] == np.int64 or df.dtypes[i] == np.float64 or df.dtypes[i] == np.int32):
        numerical_data.append(i)
    else:
        timeseries_data.append(i)
        

print("Numerical Data: ", numerical_data)
print("Categorical Data: ", categorical_data)
print("Time Series Data: ", timeseries_data)
# Space for initial data analysis ends

#Load the names of the Templates
listTxt2 = []
Templates = ["Template 1","Template 2","Template 3"]
try:
    with open("AppFinal_TemplateNames.txt") as f:
        for l in f.readlines(): # Goes line by line
            listTxt2.append(l.replace('\n',''))
except FileNotFoundError:
    Templates = ["Template 1","Template 2","Template 3"]
else:
    i = 0
    while(i < 3):
        Templates[i] = "Temp " + str(i+1) + ") " + listTxt2[i]
        i += 1

# Build the main window
btn_temp1 = tk.Button(text = Templates[0], width = 18, height = 1, bg = "white", fg = "blue")
btn_temp1.bind("<Button-1>", load_template1)
btn_temp2 = tk.Button(text = Templates[1], width = 18, height = 1, bg = "white", fg = "blue")
btn_temp2.bind("<Button-1>", load_template2)
btn_temp3 = tk.Button(text = Templates[2], width = 18, height = 1, bg = "white", fg = "blue")
btn_temp3.bind("<Button-1>", load_template3)
btn_editDataType = tk.Button(text = "Change Data Types", width = 18, height = 1, bg = "white")
btn_editDataType.bind("<Button-1>", Edit_Datatypes)
lbl_space1 = tk.Label(window, text = " ")
lbl_space2 = tk.Label(window, text = " ")
lbl_uni = tk.Label(window, text = "Univariant Chart", font = "bold")
btn_hist = tk.Button(text = "Histogram", width = 18, height = 1, bg = "white")
btn_hist.bind("<Button-1>", insert_histogram)
btn_vio = tk.Button(text = "Violin Chart", width = 18, height = 1, bg = "white")
btn_vio.bind("<Button-1>", insert_violin)
btn_pie = tk.Button(text = "Pie Chart", width = 18, height = 1, bg = "white")
btn_pie.bind("<Button-1>", insert_pie)
btn_temp1.grid(row = 0, column = 1)
btn_temp2.grid(row = 0, column = 3)
btn_temp3.grid(row = 0, column = 5)
btn_editDataType.grid(row = 0, column = 7)
lbl_uni.grid(row = 2, column = 1)
btn_hist.grid(row = 3, column = 1)
lbl_space1.grid(row = 1, column = 2)
btn_vio.grid(row = 3, column = 3)
lbl_space2.grid(row = 3, column = 4)
btn_pie.grid(row = 3, column = 5)
tk.Label(window, text = " ").grid(row = 3, column = 6)
btn_cou = tk.Button(text = "Count Chart", width = 18, height = 1, bg = "white")
btn_cou.bind("<Button-1>", insert_count)
btn_cou.grid(row = 3, column = 7)

tk.Label(window, text = " ").grid(row = 4, column = 0)
tk.Label(window, text = " ").grid(row = 4, column = 8)
lbl_bi = tk.Label(window, text = "Bivariant Charts", font = "bold")
lbl_bi.grid(row = 5, column = 1 )
btn_bar = tk.Button(window, text = "Bar Chart", width = 18, height = 1, bg = "white")
btn_bar.bind("<Button-1>", insert_bar)
btn_bar.grid(row = 6, column = 1)
btn_lin = tk.Button(window, text = "Line Chart", width = 18, height = 1, bg = "white")
btn_lin.bind("<Button-1>", insert_line)
btn_lin.grid(row = 6, column = 3)

tk.Label(window, text = " ").grid(row = 7, column = 0)
btn_submit = tk.Button(text = "Create Presentation", width=18,height=2,bg="white",fg="red")
btn_submit.bind("<Button-1>", title_slide)
btn_submit.grid(row = 8, column = 3)

tk.Label(window, text = " ").grid(row = 9, column = 0)
btn_instruct = tk.Button(window, text = "Click for Help", width = 18, height = 1, bg = "white")
btn_instruct.bind("<Button-1>", user_instructions)
btn_instruct.grid(row = 10, column = 1)


window.mainloop()
